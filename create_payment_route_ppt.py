# -*- coding: utf-8 -*-
"""
토스페이먼츠 결제경로 PPT 생성 스크립트
빌링(정기결제)용 결제경로 파일
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# PPT 생성
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 비율
prs.slide_height = Inches(7.5)

def rgb_to_hex(r, g, b):
    """RGB를 16진수로 변환"""
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)

def add_info_slide(prs, title, info_dict):
    """정보 슬라이드 추가 (표지용)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 헤더 배경 (파란색)
    header_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = rgb_to_hex(59, 130, 246)
    header_bg.line.fill.background()

    # 헤더 타이틀
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    header_frame = header_box.text_frame
    header_para = header_frame.paragraphs[0]
    header_para.text = title
    header_para.font.size = Pt(28)
    header_para.font.bold = True
    header_para.font.color.rgb = rgb_to_hex(255, 255, 255)
    header_para.alignment = PP_ALIGN.CENTER

    # 정보 박스
    y_pos = 2.0
    for key, value in info_dict.items():
        # 키 (파란색)
        key_box = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(3), Inches(0.5))
        key_frame = key_box.text_frame
        key_para = key_frame.paragraphs[0]
        key_para.text = key
        key_para.font.size = Pt(20)
        key_para.font.bold = True
        key_para.font.color.rgb = rgb_to_hex(59, 130, 246)

        # 값 (검정색)
        val_box = slide.shapes.add_textbox(Inches(5.5), Inches(y_pos), Inches(6), Inches(0.5))
        val_frame = val_box.text_frame
        val_para = val_frame.paragraphs[0]
        val_para.text = f": {value}"
        val_para.font.size = Pt(20)
        val_para.font.color.rgb = rgb_to_hex(0, 0, 0)

        y_pos += 0.7

    return slide

def add_screenshot_slide(prs, title, subtitle="", placeholder_text=""):
    """스크린샷 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 헤더 배경
    header_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = rgb_to_hex(59, 130, 246)
    header_bg.line.fill.background()

    # 헤더 타이틀
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    header_frame = header_box.text_frame
    header_para = header_frame.paragraphs[0]
    header_para.text = title
    header_para.font.size = Pt(28)
    header_para.font.bold = True
    header_para.font.color.rgb = rgb_to_hex(255, 255, 255)
    header_para.alignment = PP_ALIGN.CENTER

    # 부제목
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.333), Inches(0.5))
        sub_frame = sub_box.text_frame
        sub_para = sub_frame.paragraphs[0]
        sub_para.text = subtitle
        sub_para.font.size = Pt(14)
        sub_para.font.color.rgb = rgb_to_hex(100, 100, 100)
        sub_para.alignment = PP_ALIGN.CENTER

    # 스크린샷 플레이스홀더 (테두리 있는 박스)
    placeholder = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2), Inches(11.733), Inches(5)
    )
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = rgb_to_hex(245, 245, 245)
    placeholder.line.color.rgb = rgb_to_hex(200, 200, 200)
    placeholder.line.width = Pt(2)

    # 플레이스홀더 텍스트
    if placeholder_text:
        ph_text = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(10.333), Inches(2))
        ph_frame = ph_text.text_frame
        ph_frame.word_wrap = True
        ph_para = ph_frame.paragraphs[0]
        ph_para.text = placeholder_text
        ph_para.font.size = Pt(16)
        ph_para.font.color.rgb = rgb_to_hex(120, 120, 120)
        ph_para.alignment = PP_ALIGN.CENTER

    return slide

def add_end_slide(prs, title):
    """마지막 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = rgb_to_hex(0, 0, 0)
    title_para.alignment = PP_ALIGN.CENTER

    return slide

# ============================================
# 슬라이드 생성
# ============================================

# 1. 표지 - 가맹점 정보 기재
add_info_slide(prs, "① 가맹점 정보 기재", {
    "(1) 상호명": "머프키치",
    "(2) 사업자번호": "401-20-84647",
    "(3) URL": "https://blog-index-analyzer.vercel.app",
    "(4) Test ID": "test@blank.com",
    "(5) Test PW": "test1234!"
})

# 2. 하단정보 캡처 (사업자정보)
add_screenshot_slide(
    prs,
    "② 하단 정보 캡처",
    "필수 구성 항목: 상호명 / 대표자명 / 사업자등록번호 / 통신판매업신고번호 / 사업장주소 / 유선전화번호",
    "[홈페이지 하단 푸터 스크린샷을 여기에 붙여넣으세요]\n\nURL: https://blog-index-analyzer.vercel.app"
)

# 3. 환불규정 캡처
add_screenshot_slide(
    prs,
    "③ 환불규정 캡처 (무형상품)",
    "환불 규정을 캡처해요.",
    "[환불정책 페이지 스크린샷을 여기에 붙여넣으세요]\n\nURL: https://blog-index-analyzer.vercel.app/refund-policy"
)

# 4. 로그인/회원가입 캡처
add_screenshot_slide(
    prs,
    "④ 로그인 / 회원가입 캡처",
    "로그인 혹은 회원가입 경로를 캡처해요.",
    "[로그인 페이지 스크린샷을 여기에 붙여넣으세요]\n\nURL: https://blog-index-analyzer.vercel.app/login"
)

# 5-1. 상품선택 캡처 (요금제 페이지)
add_screenshot_slide(
    prs,
    "⑤ 상품 선택 / 구매과정 캡처 (1/3)",
    "예시) shop > 카테고리 선택 > 제품 선택 > 옵션 선택 > 구매하기 탭 등",
    "[요금제 선택 페이지 스크린샷을 여기에 붙여넣으세요]\n\nURL: https://blog-index-analyzer.vercel.app/pricing"
)

# 5-2. 상품 상세 (플랜 선택)
add_screenshot_slide(
    prs,
    "⑤ 상품 선택 / 구매과정 캡처 (2/3)",
    "상품 상세 정보 및 결제 주기 선택",
    "[플랜 상세 정보 스크린샷을 여기에 붙여넣으세요]\n\n* 서비스 제공기간(1개월/12개월)이 표시되어야 합니다"
)

# 5-3. 결제 페이지
add_screenshot_slide(
    prs,
    "⑤ 상품 선택 / 구매과정 캡처 (3/3)",
    "결제 페이지 - 주문 정보 확인",
    "[결제 페이지 스크린샷을 여기에 붙여넣으세요]\n\nURL: https://blog-index-analyzer.vercel.app/payment"
)

# 6. 카드 결제경로 캡처 (빌링 결제창)
add_screenshot_slide(
    prs,
    "⑥ 카드결제경로 캡처",
    "빌링결제의 경우 정기결제용 카드 입력창까지 캡처해야해요.",
    "[토스페이먼츠 빌링 결제창 스크린샷을 여기에 붙여넣으세요]\n\n* 결제하기 버튼 클릭 후 나타나는 토스페이먼츠 결제창"
)

# 감사합니다 슬라이드
add_end_slide(prs, "감사합니다")

# PPT 저장
output_path = r"G:\내 드라이브\developer\blog-index-analyzer\머프키치_결제경로_빌링.pptx"
prs.save(output_path)
print(f"PPT 파일 생성 완료: {output_path}")
print("\n" + "="*60)
print("스크린샷 캡처 가이드")
print("="*60)
print("""
각 슬라이드의 회색 박스를 실제 스크린샷으로 교체해주세요:

1. ② 하단 정보
   - URL: https://blog-index-analyzer.vercel.app
   - 페이지 하단의 푸터 영역 캡처
   - 사업자정보가 모두 보여야 함

2. ③ 환불규정
   - URL: https://blog-index-analyzer.vercel.app/refund-policy
   - 환불 정책 내용이 보이도록 캡처

3. ④ 로그인
   - URL: https://blog-index-analyzer.vercel.app/login
   - 로그인 폼이 보이도록 캡처

4. ⑤-1 상품선택
   - URL: https://blog-index-analyzer.vercel.app/pricing
   - 요금제 카드들이 보이도록 캡처

5. ⑤-2 플랜상세
   - URL: https://blog-index-analyzer.vercel.app/pricing
   - "정기결제 서비스 안내" 섹션이 보이도록 스크롤 후 캡처

6. ⑤-3 결제페이지
   - URL: https://blog-index-analyzer.vercel.app/payment?orderId=test&amount=29900&orderName=프로플랜&planType=pro&billingCycle=monthly
   - 결제 정보 및 약관 동의가 보이도록 캡처

7. ⑥ 결제창
   - 결제하기 버튼 클릭 후 나타나는 토스페이먼츠 팝업
   - 카드 정보 입력창 캡처

⚠️ 주의사항:
- 모든 캡처에 URL 주소창이 보여야 함
- PC 우측 하단 시간이 함께 캡처되어야 함
- 북마크바는 숨기고 캡처
""")
