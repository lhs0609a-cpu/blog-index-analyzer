# -*- coding: utf-8 -*-
"""
토스페이먼츠 결제경로 PPT 생성 스크립트
실제 스크린샷 캡처 후 PPT에 삽입
"""

import asyncio
import os
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from datetime import datetime

# 설정
BASE_URL = "https://blog-index-analyzer.vercel.app"
SCREENSHOT_DIR = r"G:\내 드라이브\developer\blog-index-analyzer\screenshots"
OUTPUT_PPT = r"G:\내 드라이브\developer\blog-index-analyzer\머프키치_결제경로_빌링_완성.pptx"

# 스크린샷 설정
PAGES_TO_CAPTURE = [
    {
        "name": "footer",
        "url": f"{BASE_URL}",
        "title": "② 하단 정보 캡처",
        "subtitle": "필수 구성 항목: 상호명 / 대표자명 / 사업자등록번호 / 통신판매업신고번호 / 사업장주소 / 유선전화번호",
        "scroll_to_bottom": True,
        "viewport": {"width": 1920, "height": 1080}
    },
    {
        "name": "refund_policy",
        "url": f"{BASE_URL}/refund-policy",
        "title": "③ 환불규정 캡처 (무형상품)",
        "subtitle": "환불 규정을 캡처해요.",
        "full_page": True,
        "viewport": {"width": 1920, "height": 1080}
    },
    {
        "name": "login",
        "url": f"{BASE_URL}/login",
        "title": "④ 로그인 / 회원가입 캡처",
        "subtitle": "로그인 혹은 회원가입 경로를 캡처해요.",
        "viewport": {"width": 1920, "height": 1080}
    },
    {
        "name": "pricing_1",
        "url": f"{BASE_URL}/pricing",
        "title": "⑤ 상품 선택 / 구매과정 캡처 (1/3)",
        "subtitle": "예시) shop > 카테고리 선택 > 제품 선택 > 옵션 선택 > 구매하기 탭 등",
        "viewport": {"width": 1920, "height": 1080}
    },
    {
        "name": "pricing_2",
        "url": f"{BASE_URL}/pricing",
        "title": "⑤ 상품 선택 / 구매과정 캡처 (2/3)",
        "subtitle": "상품 상세 정보 및 결제 주기 선택",
        "scroll_to_element": "정기결제 서비스 안내",
        "viewport": {"width": 1920, "height": 1080}
    },
    {
        "name": "payment",
        "url": f"{BASE_URL}/payment?orderId=TEST001&amount=29900&orderName=프로플랜(월간)&planType=pro&billingCycle=monthly",
        "title": "⑤ 상품 선택 / 구매과정 캡처 (3/3)",
        "subtitle": "결제 페이지 - 주문 정보 확인",
        "viewport": {"width": 1920, "height": 1080}
    }
]


async def capture_screenshots():
    """웹 페이지 스크린샷 캡처"""
    os.makedirs(SCREENSHOT_DIR, exist_ok=True)
    screenshots = {}

    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)

        for page_config in PAGES_TO_CAPTURE:
            context = await browser.new_context(
                viewport=page_config.get("viewport", {"width": 1920, "height": 1080}),
                locale="ko-KR"
            )
            page = await context.new_page()

            print(f"캡처 중: {page_config['name']} - {page_config['url']}")

            await page.goto(page_config["url"], wait_until="networkidle")
            await asyncio.sleep(2)  # 페이지 로딩 대기

            # 스크롤 처리
            if page_config.get("scroll_to_bottom"):
                await page.evaluate("window.scrollTo(0, document.body.scrollHeight)")
                await asyncio.sleep(1)

            if page_config.get("scroll_to_element"):
                try:
                    element = await page.query_selector(f"text={page_config['scroll_to_element']}")
                    if element:
                        await element.scroll_into_view_if_needed()
                        await asyncio.sleep(1)
                except:
                    pass

            # 스크린샷 저장
            screenshot_path = os.path.join(SCREENSHOT_DIR, f"{page_config['name']}.png")

            if page_config.get("full_page"):
                await page.screenshot(path=screenshot_path, full_page=True)
            else:
                await page.screenshot(path=screenshot_path)

            screenshots[page_config["name"]] = {
                "path": screenshot_path,
                "title": page_config["title"],
                "subtitle": page_config["subtitle"]
            }

            await context.close()

        await browser.close()

    return screenshots


def create_ppt_with_screenshots(screenshots):
    """스크린샷을 포함한 PPT 생성"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_info_slide(title, info_dict):
        """정보 슬라이드 추가"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 헤더 배경
        header_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
        )
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = RGBColor(59, 130, 246)
        header_bg.line.fill.background()

        # 헤더 타이틀
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
        header_para = header_box.text_frame.paragraphs[0]
        header_para.text = title
        header_para.font.size = Pt(28)
        header_para.font.bold = True
        header_para.font.color.rgb = RGBColor(255, 255, 255)
        header_para.alignment = PP_ALIGN.CENTER

        # 정보 표시
        y_pos = 2.0
        for key, value in info_dict.items():
            key_box = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(3), Inches(0.5))
            key_para = key_box.text_frame.paragraphs[0]
            key_para.text = key
            key_para.font.size = Pt(20)
            key_para.font.bold = True
            key_para.font.color.rgb = RGBColor(59, 130, 246)

            val_box = slide.shapes.add_textbox(Inches(5.5), Inches(y_pos), Inches(6), Inches(0.5))
            val_para = val_box.text_frame.paragraphs[0]
            val_para.text = f": {value}"
            val_para.font.size = Pt(20)
            val_para.font.color.rgb = RGBColor(0, 0, 0)

            y_pos += 0.7

    def add_screenshot_slide(title, subtitle, image_path):
        """스크린샷 슬라이드 추가"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 헤더 배경
        header_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0)
        )
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = RGBColor(59, 130, 246)
        header_bg.line.fill.background()

        # 헤더 타이틀
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.5))
        header_para = header_box.text_frame.paragraphs[0]
        header_para.text = title
        header_para.font.size = Pt(24)
        header_para.font.bold = True
        header_para.font.color.rgb = RGBColor(255, 255, 255)
        header_para.alignment = PP_ALIGN.CENTER

        # 부제목
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.4))
            sub_para = sub_box.text_frame.paragraphs[0]
            sub_para.text = subtitle
            sub_para.font.size = Pt(12)
            sub_para.font.color.rgb = RGBColor(100, 100, 100)
            sub_para.alignment = PP_ALIGN.CENTER

        # 스크린샷 이미지 삽입
        if os.path.exists(image_path):
            # 이미지 크기 계산 (슬라이드에 맞게)
            left = Inches(0.3)
            top = Inches(1.5)
            width = Inches(12.733)
            height = Inches(5.8)

            slide.shapes.add_picture(image_path, left, top, width=width, height=height)

    def add_placeholder_slide(title, subtitle, placeholder_text):
        """플레이스홀더 슬라이드 (스크린샷이 없을 때)"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 헤더 배경
        header_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0)
        )
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = RGBColor(59, 130, 246)
        header_bg.line.fill.background()

        # 헤더 타이틀
        header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.5))
        header_para = header_box.text_frame.paragraphs[0]
        header_para.text = title
        header_para.font.size = Pt(24)
        header_para.font.bold = True
        header_para.font.color.rgb = RGBColor(255, 255, 255)
        header_para.alignment = PP_ALIGN.CENTER

        # 부제목
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.4))
            sub_para = sub_box.text_frame.paragraphs[0]
            sub_para.text = subtitle
            sub_para.font.size = Pt(12)
            sub_para.font.color.rgb = RGBColor(100, 100, 100)
            sub_para.alignment = PP_ALIGN.CENTER

        # 플레이스홀더 박스
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.6)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(245, 245, 245)
        placeholder.line.color.rgb = RGBColor(200, 200, 200)
        placeholder.line.width = Pt(2)

        # 플레이스홀더 텍스트
        ph_text = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(10.333), Inches(2))
        ph_frame = ph_text.text_frame
        ph_frame.word_wrap = True
        ph_para = ph_frame.paragraphs[0]
        ph_para.text = placeholder_text
        ph_para.font.size = Pt(16)
        ph_para.font.color.rgb = RGBColor(120, 120, 120)
        ph_para.alignment = PP_ALIGN.CENTER

    def add_end_slide(title):
        """마지막 슬라이드"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
        title_para = title_box.text_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 0, 0)
        title_para.alignment = PP_ALIGN.CENTER

    # ============================================
    # 슬라이드 생성
    # ============================================

    # 1. 가맹점 정보
    add_info_slide("① 가맹점 정보 기재", {
        "(1) 상호명": "머프키치",
        "(2) 사업자번호": "401-20-84647",
        "(3) URL": "https://blog-index-analyzer.vercel.app",
        "(4) Test ID": "test@blank.com",
        "(5) Test PW": "test1234!"
    })

    # 2. 하단 정보 (푸터)
    if "footer" in screenshots:
        add_screenshot_slide(
            screenshots["footer"]["title"],
            screenshots["footer"]["subtitle"],
            screenshots["footer"]["path"]
        )

    # 3. 환불규정
    if "refund_policy" in screenshots:
        add_screenshot_slide(
            screenshots["refund_policy"]["title"],
            screenshots["refund_policy"]["subtitle"],
            screenshots["refund_policy"]["path"]
        )

    # 4. 로그인
    if "login" in screenshots:
        add_screenshot_slide(
            screenshots["login"]["title"],
            screenshots["login"]["subtitle"],
            screenshots["login"]["path"]
        )

    # 5-1. 상품선택 (요금제)
    if "pricing_1" in screenshots:
        add_screenshot_slide(
            screenshots["pricing_1"]["title"],
            screenshots["pricing_1"]["subtitle"],
            screenshots["pricing_1"]["path"]
        )

    # 5-2. 플랜 상세
    if "pricing_2" in screenshots:
        add_screenshot_slide(
            screenshots["pricing_2"]["title"],
            screenshots["pricing_2"]["subtitle"],
            screenshots["pricing_2"]["path"]
        )

    # 5-3. 결제 페이지
    if "payment" in screenshots:
        add_screenshot_slide(
            screenshots["payment"]["title"],
            screenshots["payment"]["subtitle"],
            screenshots["payment"]["path"]
        )

    # 6. 카드 결제경로 (토스페이먼츠 팝업) - 수동 캡처 필요
    add_placeholder_slide(
        "⑥ 카드결제경로 캡처",
        "빌링결제의 경우 정기결제용 카드 입력창까지 캡처해야해요.",
        "[토스페이먼츠 빌링 결제창 스크린샷을 여기에 붙여넣으세요]\n\n* 결제하기 버튼 클릭 후 나타나는 토스페이먼츠 결제창\n* 이 화면은 실제 결제 버튼을 클릭해서 수동 캡처가 필요합니다"
    )

    # 감사합니다
    add_end_slide("감사합니다")

    # PPT 저장
    prs.save(OUTPUT_PPT)
    print(f"\nPPT 파일 생성 완료: {OUTPUT_PPT}")


async def main():
    print("=" * 60)
    print("토스페이먼츠 결제경로 PPT 자동 생성")
    print(f"시작 시간: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    print("=" * 60)

    print("\n[1/2] 스크린샷 캡처 중...")
    screenshots = await capture_screenshots()
    print(f"캡처 완료: {len(screenshots)}개 페이지")

    print("\n[2/2] PPT 생성 중...")
    create_ppt_with_screenshots(screenshots)

    print("\n" + "=" * 60)
    print("완료!")
    print("=" * 60)
    print(f"\n생성된 파일: {OUTPUT_PPT}")
    print(f"스크린샷 폴더: {SCREENSHOT_DIR}")
    print("\n주의: ⑥ 카드결제경로 슬라이드는 토스페이먼츠 팝업을 수동 캡처해서 붙여넣어야 합니다.")


if __name__ == "__main__":
    asyncio.run(main())
